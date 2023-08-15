from pptx import Presentation

prs = Presentation('sample.pptx')

text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
    text_runs.append("\n")
        

with open('sample_text.txt', 'w') as textfile:
    for line in text_runs:
        print(line)
        try:
            textfile.write(str(line) + "\n")
        except:
            pass
    