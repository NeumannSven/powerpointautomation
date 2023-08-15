from pptx import Presentation

prs = Presentation('prs005.pptx')

text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_table:
            for i in shape.table.rows:
                tabellen_text = ''
                for j in i.cells:
                    tabellen_text += j.text + "\t"
                    
                text_runs.append(tabellen_text)
                
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
    text_runs.append("\n")


prs.slides[2].shapes[1].table.rows[1].cells[0].text = "42"

prs.save("prs005.pptx")    
    
with open('prs005_text.txt', 'w') as textfile:
    for line in text_runs:
        print(line)
        try:
            textfile.write(str(line) + "\n")
        except:
            pass
    