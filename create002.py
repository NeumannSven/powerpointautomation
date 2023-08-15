from pptx import Presentation

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "pySpaceBremen"
subtitle.text = "PowerPoint mit Python automatisieren"



bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Agenda'

tf = body_shape.text_frame
tf.text = 'Installation'

p = tf.add_paragraph()
p.text = 'Python'
p.level = 1

p = tf.add_paragraph()
p.text = 'python.org'
p.level = 2

p = tf.add_paragraph()
p.text = 'PyPI'
p.level = 1

p = tf.add_paragraph()
p.text = 'pypi.org'
p.level = 2

p = tf.add_paragraph()
p.text = 'pip install python-pptx'
p.level = 2

prs.save('prs002.pptx')